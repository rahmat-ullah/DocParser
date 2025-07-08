"""
PPTX Parser using python-pptx.
Extracts text, images, and other content from PowerPoint files.
"""

import base64
from typing import Optional, AsyncGenerator
from pathlib import Path
from pptx import Presentation

from .base_parser import BaseParser, ParseError
from .ast_models import DocumentAST, TextBlock, ImageBlock, BlockType, ParseProgress


class PPTXParser(BaseParser):
    """Parser for PPTX documents using python-pptx."""

    def supports_file(self, file_path: Path) -> bool:
        """Check if file is a PPTX."""
        return file_path.suffix.lower() == '.pptx'

    async def parse(
        self, file_path: Path, progress_callback: Optional[AsyncGenerator[ParseProgress, None]] = None
    ) -> DocumentAST:
        """Parse PPTX document and extract content."""
        try:
            await self._emit_progress(progress_callback, "initialization", 0.0, "Opening PPTX document")

            presentation = Presentation(file_path)
            ast = DocumentAST(metadata={"format": "PPTX", "slides": len(presentation.slides)})

            total_slides = len(presentation.slides)
            
            for i, slide in enumerate(presentation.slides):
                await self._emit_progress(
                    progress_callback, 
                    "parsing_slides", 
                    i / total_slides, 
                    f"Processing slide {i + 1} of {total_slides}"
                )
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                # Determine if it's a title or content
                                block_type = BlockType.HEADING if shape == slide.shapes.title else BlockType.PARAGRAPH
                                level = 1 if block_type == BlockType.HEADING else None
                                
                                text_block = TextBlock(
                                    type=block_type,
                                    content=paragraph.text.strip(),
                                    level=level
                                )
                                ast.textBlocks.append(text_block)

                # Extract images
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_data = image.blob
                        image_base64 = base64.b64encode(image_data).decode()
                        image_block = ImageBlock(
                            data=image_base64,
                            format=image.ext.upper()
                        )
                        ast.images.append(image_block)

            await self._emit_progress(progress_callback, "completion", 1.0, "PPTX parsing completed")
            return ast

        except Exception as e:
            raise ParseError(f"Failed to parse PPTX: {str(e)}", file_path, e)
