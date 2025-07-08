import pytest
import asyncio
from pathlib import Path
from pptx import Presentation

from app.parsers.pptx_parser import PPTXParser

@pytest.fixture
def pptx_parser():
    return PPTXParser()

@pytest.fixture
def pptx_file(tmp_path):
    # Create a simple PPTX file with content
    prs = Presentation()
    
    # Add a slide with title and content
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Test Presentation"
    
    content = slide.placeholders[1]
    content.text = "This is test content for the presentation."
    
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path

@pytest.mark.asyncio
async def test_pptx_parser_supports_file(pptx_parser, pptx_file):
    assert pptx_parser.supports_file(pptx_file) == True
    
@pytest.mark.asyncio
async def test_pptx_parser_parse(pptx_parser, pptx_file):
    result = await pptx_parser.parse_to_dict(pptx_file)
    assert isinstance(result, dict)
    assert "metadata" in result
    assert result["metadata"]["format"] == "PPTX"
    assert "textBlocks" in result
    assert len(result["textBlocks"]) > 0
    
    # Check if title was detected as heading
    headings = [block for block in result["textBlocks"] if block["type"] == "heading"]
    assert len(headings) > 0
    
    # Check if content was detected
    paragraphs = [block for block in result["textBlocks"] if block["type"] == "paragraph"]
    assert len(paragraphs) > 0
